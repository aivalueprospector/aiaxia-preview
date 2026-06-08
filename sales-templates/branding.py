#!/usr/bin/env python3
"""
Generate branded pandoc reference docs (.docx and .pptx) for each brand.

These are CONSUMED by pandoc via --reference-doc: pandoc copies their named
styles / slide layouts into the rendered output. Run once before build.py
(build.py calls this automatically if a reference doc is missing).

Output: reference-docs/{brand}.docx  and  reference-docs/{brand}.pptx
"""
import json
import os
from pathlib import Path

from docx import Document
from docx.enum.style import WD_STYLE_TYPE
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from docx.shared import Pt, Inches, RGBColor, Emu

HERE = Path(__file__).resolve().parent
REPO = HERE.parent  # aiaxia-preview/
REFDIR = HERE / "reference-docs"
BODY_FONT = "Inter"


def _rgb(hexstr):
    return RGBColor.from_string(hexstr)


def _set_cell_shading(paragraph, hex_fill):
    """Apply paragraph background shading (w:shd) — used for callout boxes."""
    pPr = paragraph._p.get_or_add_pPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), hex_fill)
    pPr.append(shd)


def _add_page_number_field(paragraph):
    run = paragraph.add_run()
    fldBegin = OxmlElement("w:fldChar")
    fldBegin.set(qn("w:fldCharType"), "begin")
    instr = OxmlElement("w:instrText")
    instr.set(qn("xml:space"), "preserve")
    instr.text = "PAGE"
    fldEnd = OxmlElement("w:fldChar")
    fldEnd.set(qn("w:fldCharType"), "end")
    run._r.append(fldBegin)
    run._r.append(instr)
    run._r.append(fldEnd)


def build_docx(brand_key, b):
    doc = Document()

    # --- page geometry: Letter, 0.6in margins ---
    sec = doc.sections[0]
    sec.page_height = Inches(11)
    sec.page_width = Inches(8.5)
    for m in ("top_margin", "bottom_margin", "left_margin", "right_margin"):
        setattr(sec, m, Inches(0.7) if "top" in m or "bottom" in m else Inches(0.8))

    # --- base / Normal ---
    normal = doc.styles["Normal"]
    normal.font.name = BODY_FONT
    normal.font.size = Pt(10.5)
    normal.font.color.rgb = _rgb("1A2230")
    normal.paragraph_format.space_after = Pt(6)
    normal.paragraph_format.line_spacing = 1.25

    # --- headings ---
    def style_heading(name, size, color_hex, before, after):
        st = doc.styles[name]
        st.font.name = BODY_FONT
        st.font.size = Pt(size)
        st.font.bold = True
        st.font.color.rgb = _rgb(color_hex)
        st.paragraph_format.space_before = Pt(before)
        st.paragraph_format.space_after = Pt(after)

    style_heading("Title", 26, b["accent_dark"], 0, 6)
    if "Subtitle" in [s.name for s in doc.styles]:
        sub = doc.styles["Subtitle"]
        sub.font.name = BODY_FONT
        sub.font.size = Pt(12)
        sub.font.color.rgb = _rgb("45526A")
    style_heading("Heading 1", 19, b["accent_dark"], 16, 6)
    style_heading("Heading 2", 15, b["ink"], 12, 4)
    style_heading("Heading 3", 12.5, b["accent_dark"], 10, 2)

    # --- Placeholder character style (highlighted {{token}}) ---
    ph = doc.styles.add_style("Placeholder", WD_STYLE_TYPE.CHARACTER)
    ph.font.name = BODY_FONT
    ph.font.bold = True
    ph.font.color.rgb = _rgb("7A4B00")
    ph.font.highlight_color = WD_COLOR_INDEX.YELLOW

    # --- Block Text (pandoc blockquote → HOW-TO callout) with soft fill ---
    if "Block Text" in [s.name for s in doc.styles]:
        bt = doc.styles["Block Text"]
    else:
        bt = doc.styles.add_style("Block Text", WD_STYLE_TYPE.PARAGRAPH)
    bt.font.name = BODY_FONT
    bt.font.size = Pt(10)
    bt.font.color.rgb = _rgb("3A2F10")
    bt.paragraph_format.left_indent = Inches(0.15)
    bt.paragraph_format.right_indent = Inches(0.15)
    bt.paragraph_format.space_before = Pt(6)
    bt.paragraph_format.space_after = Pt(6)
    # soft brand-tint fill + left accent border on the callout box
    bt_pPr = bt._element.get_or_add_pPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), b["soft"])
    bt_pPr.append(shd)
    pbdr = OxmlElement("w:pBdr")
    left = OxmlElement("w:left")
    left.set(qn("w:val"), "single")
    left.set(qn("w:sz"), "18")
    left.set(qn("w:space"), "8")
    left.set(qn("w:color"), b["accent"])
    pbdr.append(left)
    bt_pPr.append(pbdr)

    # --- header: logo + brand line ---
    header = sec.header
    hp = header.paragraphs[0]
    hp.alignment = WD_ALIGN_PARAGRAPH.LEFT
    logo_path = REPO / b["logo"]
    if logo_path.exists():
        try:
            hp.add_run().add_picture(str(logo_path), height=Inches(0.28))
        except Exception:
            pass
    run = hp.add_run("   " + b["brand_line"])
    run.font.name = BODY_FONT
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = _rgb(b["accent_dark"])

    # --- footer: brand domain + page number ---
    footer = sec.footer
    fp = footer.paragraphs[0]
    fr = fp.add_run(b["display"] + " · " + b["domain"] + "    ·    Sales template · page ")
    fr.font.name = BODY_FONT
    fr.font.size = Pt(8)
    fr.font.color.rgb = _rgb("6B7A8A")
    _add_page_number_field(fp)

    # --- table styling default (light grid handled by pandoc 'Table' style) ---

    REFDIR.mkdir(exist_ok=True)
    out = REFDIR / f"{brand_key}.docx"
    doc.save(str(out))
    return out


def build_pptx(brand_key, b):
    from pptx import Presentation
    from pptx.util import Pt as PPt, Inches as PInches
    from pptx.dml.color import RGBColor as PRGB

    prs = Presentation()  # default 4:3? -> set 16:9
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)

    accent_dark = PRGB.from_string(b["accent_dark"])
    logo_path = REPO / b["logo"]
    master = prs.slide_masters[0]

    # Recolor titles -> accent_dark and set body font Inter across master + layouts.
    def restyle(container, color=True):
        for ph in container.placeholders:
            try:
                is_title = ph.placeholder_format.idx == 0 or "Title" in (ph.name or "")
                for para in ph.text_frame.paragraphs:
                    para.font.name = BODY_FONT
                    if color and is_title:
                        para.font.color.rgb = accent_dark
                        para.font.bold = True
                    for run in para.runs:
                        run.font.name = BODY_FONT
            except Exception:
                pass

    restyle(master)
    for layout in master.slide_layouts:
        restyle(layout)
        # logo top-right on each layout (LayoutShapes supports add_picture)
        if logo_path.exists():
            try:
                layout.shapes.add_picture(
                    str(logo_path), PInches(12.35), PInches(0.18), height=PInches(0.55)
                )
            except Exception:
                pass

    REFDIR.mkdir(exist_ok=True)
    out = REFDIR / f"{brand_key}.pptx"
    prs.save(str(out))
    return out


def main():
    brands = json.loads((HERE / "brands.json").read_text())
    for key, b in brands.items():
        d = build_docx(key, b)
        p = build_pptx(key, b)
        print(f"  {key}: {d.name}, {p.name}")


if __name__ == "__main__":
    print("Building reference docs...")
    main()
    print("Done.")
